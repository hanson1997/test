#!/usr/bin/env python3
"""EDFIP presentation — PowerPoint for the Emeraid meeting."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from lxml import etree
from copy import deepcopy

BG = RGBColor(0x0B, 0x0F, 0x0D)
INK = RGBColor(0xF4, 0xF0, 0xE6)
MUTED = RGBColor(0xA8, 0xB3, 0xAA)
GOLD = RGBColor(0xC4, 0xA0, 0x56)
PANEL = RGBColor(0x12, 0x18, 0x16)
GREEN = RGBColor(0x3F, 0x8F, 0x6E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def paint_bg(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()
    # send to back
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_text(slide, l, t, w, h, text, size=18, color=INK, bold=False, font="Georgia", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font
    return box


def kicker_title(slide, kicker, title, num):
    paint_bg(slide)
    add_text(slide, 0.7, 0.35, 11, 0.35, kicker.upper(), 12, GOLD, True, "Calibri")
    add_text(slide, 0.7, 0.75, 12, 1.6, title, 32, INK, False, "Georgia")
    add_text(slide, 0.7, 7.05, 8, 0.25, "Dexta Synergy Services  ·  EDFIP", 11, MUTED, False, "Calibri")
    add_text(slide, 11.4, 7.05, 1.4, 0.25, f"{num:02d} / 15", 11, MUTED, False, "Calibri", PP_ALIGN.RIGHT)


def bullets(slide, items, top=2.5, size=18):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(12), Inches(4.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        run = p.add_run()
        run.text = "•  " + item
        run.font.size = Pt(size)
        run.font.color.rgb = INK
        run.font.name = "Calibri"


# 1
s = prs.slides.add_slide(blank)
paint_bg(s)
add_text(s, 0.7, 0.4, 12, 0.35, "HOW WE READ THE ASSIGNMENT", 13, GOLD, True, "Calibri")
add_text(s, 0.7, 0.9, 12, 1.5, "We understand the product\nEmeraid is building.", 34, INK, False, "Georgia")
add_text(s, 0.7, 2.6, 11.5, 1.8, "Not a basic microfinance application. A commercialisable, multi-tenant Digital Financial Inclusion Operating System — EDFIP — that Emeraid will own, control, brand and sell to third-party institutions.", 18, MUTED, False, "Calibri")
add_text(s, 0.7, 4.6, 12, 1.6, "Microfinance institutions   ·   Cooperatives   ·   VSLA networks\nNGOs and programmes   ·   Solar / PayGo operators   ·   Agency banking networks", 16, INK, False, "Calibri")
add_text(s, 0.7, 6.5, 10, 0.35, "Every architectural and commercial choice in our proposal follows from that reality.", 15, GOLD, False, "Calibri")
add_text(s, 0.7, 7.05, 8, 0.25, "Dexta Synergy Services  ·  EDFIP", 11, MUTED, False, "Calibri")
add_text(s, 11.4, 7.05, 1.4, 0.25, "01 / 15", 11, MUTED, False, "Calibri", PP_ALIGN.RIGHT)

# 2
s = prs.slides.add_slide(blank)
kicker_title(s, "What this engagement must protect", "Four commitments that matter\nto Emeraid.", 2)
bullets(s, [
    "Full CRM in the first production release — the complete relationship, not a registration screen.",
    "Green Asset Finance and PayGo in the first release, with an Early Operational Release by Month 6.",
    "Emeraid’s ability to operate, modify, commercialise and maintain the platform independently.",
    "A realistic 12-month delivery plan with testing and acceptance evidence.",
], 2.6, 17)

# 3
s = prs.slides.add_slide(blank)
kicker_title(s, "First-release scope", "The platform Emeraid specified.", 3)
add_text(s, 0.7, 2.45, 4, 4.3, "OPERATE\n\nMulti-tenant SaaS Emeraid can sell\nClients, members and KYC\nFull CRM\nRoles, maker-checker and audit", 14, INK, False, "Calibri")
add_text(s, 4.8, 2.45, 4, 4.3, "LEND AND SAVE\n\nLoan management\nSavings, including esusu / ajo\nAccounting and general ledger\nCooperatives and VSLA groups", 14, INK, False, "Calibri")
add_text(s, 8.9, 2.45, 4, 4.3, "DIFFERENTIATE\n\nGreen Asset Finance and PayGo\nDevice-linked credit — not cash-only\nAgency banking and tellering\nDigital payments and collections", 14, INK, False, "Calibri")
add_text(s, 0.7, 6.35, 12, 0.55, "Reach: offline field app · customer web + Android · SMS/email    Prove: dashboards, donor reporting, partner APIs, first-tenant migration", 13, MUTED, False, "Calibri")

# 4
s = prs.slides.add_slide(blank)
kicker_title(s, "Multi-tenant Software as a Service", "One control plane. Emeraid onboards.\nModules switch on.", 4)
add_text(s, 0.7, 2.7, 3.8, 3.4, "FIRST RELEASE\n\nOne platform on Emeraid’s server.\nInstitutions get a login.\nNo separate install per client.", 16, INK, False, "Calibri")
add_text(s, 4.8, 2.7, 3.8, 3.4, "WHAT THEY BUY FROM YOU\n\nA module pack — CRM only,\nVSLA only, full MFI, or PayGo\noperator. Super-admin turns it on.", 16, INK, False, "Calibri")
add_text(s, 8.9, 2.7, 3.8, 3.4, "IF A LICENCE LAPSES\n\nThe institution can be suspended.\nAccess locked. Data retained.\nRestored when the licence is current.", 16, INK, False, "Calibri")

# 5
s = prs.slides.add_slide(blank)
kicker_title(s, "How Emeraid licences the product", "Same platform. Three different products.", 5)
add_text(s, 0.7, 2.5, 3.9, 4.2, "HOPE MICROFINANCE\nFull MFI suite\n\nON  CRM\nON  Loans + savings + GL\nON  Agency + payments\nON  PayGo\nOFF VSLA", 15, INK, False, "Calibri")
add_text(s, 4.8, 2.5, 3.9, 4.2, "RIVERS VSLA NETWORK\nProgramme pack\n\nON  CRM\nON  VSLA + share-out\nON  Donor reports\nOFF Loans / agency / PayGo", 15, INK, False, "Calibri")
add_text(s, 8.9, 2.5, 3.9, 4.2, "GREENLIGHT SOLAR\nGreen operator\n\nON  CRM + clients\nON  Loans + GL\nON  PayGo + SMS\nOFF VSLA / cooperative", 15, INK, False, "Calibri")

# 6
s = prs.slides.add_slide(blank)
kicker_title(s, "Architecture", "One runtime. One database. One ledger.", 6)
bullets(s, [
    "Odoo Community 19 — spine: CRM, general ledger, users, portal.",
    "Emeraid-owned modules — loans, savings, VSLA, cooperatives, PayGo, agency. Purpose-built for EDFIP.",
    "FastAPI inside Odoo — public API, webhooks, mobile sync. Meets the TOR security and documentation standard.",
    "Flutter — field Android app + customer Android app. Encrypted offline store.",
    "Nothing writes to PostgreSQL except through Odoo. Isolation is enforced once.",
], 2.5, 16)

# 7
s = prs.slides.add_slide(blank)
kicker_title(s, "Implementation approach", "What the foundation provides —\nand what we build.", 7)
add_text(s, 0.7, 2.5, 5.8, 4.2, "ODOO ALREADY GIVES\n\n• CRM pipeline, activities, campaigns\n• Double-entry ledger (posted journals cannot be deleted)\n• Users, record rules, customer web portal\n• Payment-provider hook, jobs", 16, INK, False, "Calibri")
add_text(s, 6.8, 2.5, 5.8, 4.2, "WE BUILD AS EMERAID MODULES\n\n• Loan engine, PAR, restructuring\n• Savings, esusu/ajo, fixed deposits\n• VSLA meetings and share-out\n• PayGo tokens and OEM connectors\n• Agency float, teller, offline sync", 16, INK, False, "Calibri")

# 8
s = prs.slides.add_slide(blank)
kicker_title(s, "Odoo  ·  FastAPI  ·  Flutter", "Which tool carries which problem.", 8)
bullets(s, [
    "Sell to many institutions, turn modules on/off — Odoo + custom tenant module",
    "Leads and campaigns — Odoo CRM   |   Complaints + client 360 — custom Odoo",
    "Official books — Odoo accounting   |   Loans, savings, VSLA, PayGo, float — custom Odoo",
    "Village work with no network — Flutter",
    "Paystack, OEMs, USSD, partners — FastAPI",
    "Customer on the web — Odoo portal   |   Customer on a phone — Flutter",
], 2.5, 16)

# 9
s = prs.slides.add_slide(blank)
kicker_title(s, "Green Asset Finance  ·  Pay-As-You-Go", "Repayment becomes light.", 9)
add_text(s, 0.7, 2.45, 12, 0.6, "Customer pays  →  Days of use earned  →  OEM connector  →  Token  →  SMS / agent / app", 16, GOLD, False, "Calibri")
bullets(s, [
    "First release: asset registry, device IDs, one live OEM, eligibility engine, audit trail, ownership transfer.",
    "Month 6 early operational release: registry, asset-linked loan, first connector, token request, delivery and audit — ahead of full go-live.",
    "Connector built against a stub so an OEM delay cannot stall engineering. New OEM = new small module, no change to core.",
], 3.2, 16)

# 10
s = prs.slides.add_slide(blank)
kicker_title(s, "Field operations", "Works in the field with no network.", 10)
bullets(s, [
    "Loan officer, CRM officer, VSLA facilitator, agent — full workflow with no network.",
    "Money is append-only. No edited balances on the phone. No silent duplicates.",
    "Sync slice starts in Month 2–3, not at the end. Tests run on every merge.",
    "Customer self-service: Odoo web portal AND Flutter Android app, both in first release.",
], 2.6, 18)

# 11
s = prs.slides.add_slide(blank)
kicker_title(s, "Twelve months to acceptance", "PayGo is on the clock, not at the end.", 11)
bullets(s, [
    "Month 1 — architecture, screens, licence register, server check, OEM checklist",
    "Months 2–3 — tenants, clients, roles, ledger spine, first sync slice",
    "Months 4–6 — loans, savings, teller, PayGo early operational release",
    "Months 6–8 — cooperatives, VSLA, full CRM",
    "Months 8–11 — agency, payments, APIs, mobile, security and UAT",
    "Month 12 — migration, go-live, training, handover · then 30 days hypercare + 12 months warranty",
], 2.5, 16)

# 12
s = prs.slides.add_slide(blank)
kicker_title(s, "Who does the work", "One technical authority. A named delivery team.", 12)
add_text(s, 0.7, 2.5, 6, 4, "HANSON EYUREN\n\nLead Developer. Architecture, financial core, API, PayGo, payments, USSD. Single point of accountability. Production Odoo today. Prior: bank USSD portals, compliance / NFIU-style controls, loan portal.", 16, INK, False, "Calibri")
add_text(s, 7.0, 2.5, 5.6, 4, "DEXTA SYNERGY SERVICES\n\nEyuren Alison — Flutter field + customer apps, offline sync, QA.\n\nDaniel Azu — project manager, 100% on EDFIP.\n\nRowland Lawson — security review support.\n\nRC 7377341", 16, INK, False, "Calibri")

# 13
s = prs.slides.add_slide(blank)
kicker_title(s, "Financial proposal", "₦70 million  ·  mandatory scope, excluding VAT", 13)
add_text(s, 0.7, 2.5, 12, 0.5, "₦75.25 million including VAT.  840 person-days.  12 months.  Warranty included.", 18, MUTED, False, "Calibri")
add_text(s, 0.7, 3.3, 3.8, 2.8, "INCEPTION\n\n10% at accepted D1.\nBelow the 15% ceiling.", 16, INK, False, "Calibri")
add_text(s, 4.8, 3.3, 3.8, 2.8, "PAYGO AT MONTH 6\n\n12% on accepted Deliverable 5.\nPayment follows a working PayGo release.", 16, INK, False, "Calibri")
add_text(s, 8.9, 3.3, 3.8, 2.8, "NO HOSTING IN THE PRICE\n\nEmeraid’s server.\nWe install, harden, hand over.", 16, INK, False, "Calibri")

# 14
s = prs.slides.add_slide(blank)
kicker_title(s, "Ownership and handover", "Emeraid owns the platform —\nand can run it independently.", 14)
bullets(s, [
    "Free to operate, modify, extend, sell, and onboard tenants.",
    "Source repository is Emeraid’s from day one. Independent maintenance after handover.",
    "No Odoo Enterprise. No paid App Store. Licence register regenerated on every commit.",
    "Acceptance includes a live restore and a deploy by Emeraid staff.",
], 2.6, 18)

# 15
s = prs.slides.add_slide(blank)
paint_bg(s)
add_text(s, 0.7, 0.45, 12, 0.35, "THANK YOU", 13, GOLD, True, "Calibri")
add_text(s, 0.7, 2.0, 12, 1.6, "We welcome your\nquestions.", 40, INK, False, "Georgia")
add_text(s, 0.7, 4.3, 11, 1.2, "Architecture, PayGo, CRM, delivery capacity, commercial terms,\nand any collaboration model you wish to explore.", 20, MUTED, False, "Calibri")
add_text(s, 0.7, 6.2, 11, 0.4, "Hanson Eyuren  ·  08106248715  ·  Dexta Synergy Services", 18, GOLD, False, "Calibri")
add_text(s, 11.4, 7.05, 1.4, 0.25, "15 / 15", 11, MUTED, False, "Calibri", PP_ALIGN.RIGHT)

out = "/workspace/edfip-screening/EDFIP_Screening_Dexta.pptx"
prs.save(out)
print("wrote", out)
