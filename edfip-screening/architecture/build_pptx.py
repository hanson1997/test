#!/usr/bin/env python3
"""Client-facing EDFIP architecture walkthrough for Emeraid."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BG = RGBColor(0x0B, 0x0F, 0x0D)
INK = RGBColor(0xF4, 0xF0, 0xE6)
MUTED = RGBColor(0xA8, 0xB3, 0xAA)
GOLD = RGBColor(0xC4, 0xA0, 0x56)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]
TOTAL = 16


def paint_bg(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_text(slide, l, t, w, h, text, size=18, color=INK, bold=False, font="Georgia", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font
    return box


def footer(slide, num):
    add_text(slide, 0.7, 7.05, 9, 0.25, "Dexta Synergy Services  ·  Confidential", 11, MUTED, False, "Calibri")
    add_text(slide, 11.2, 7.05, 1.6, 0.25, f"{num:02d} / {TOTAL:02d}", 11, MUTED, False, "Calibri", PP_ALIGN.RIGHT)


def kicker_title(slide, kicker, title, num):
    paint_bg(slide)
    add_text(slide, 0.7, 0.35, 12, 0.35, kicker.upper(), 12, GOLD, True, "Calibri")
    add_text(slide, 0.7, 0.75, 12, 1.55, title, 30, INK, False, "Georgia")
    footer(slide, num)


# 1
s = prs.slides.add_slide(blank)
paint_bg(s)
add_text(s, 0.7, 0.4, 12, 0.35, "PROPOSED ARCHITECTURE", 13, GOLD, True, "Calibri")
add_text(s, 0.7, 1.0, 12, 2.2, "One platform Emeraid hosts.\nInstitutions are onboarded\nand licensed.", 34, INK, False, "Georgia")
add_text(s, 0.7, 3.6, 11.5, 1.8, "This document sets out how EDFIP operates as multi-tenant software-as-a-service that Emeraid will host, brand, and sell.", 18, MUTED, False, "Calibri")
footer(s, 1)

# 2
s = prs.slides.add_slide(blank)
kicker_title(s, "Software as a service", "Emeraid does not install EDFIP\nat each institution.", 2)
add_text(s, 0.7, 2.55, 2.8, 3.2, "1. HOST\n\nOne platform on Emeraid’s environment.", 15, INK, False, "Calibri")
add_text(s, 3.7, 2.55, 2.8, 3.2, "2. ONBOARD\n\nSuper Administrator creates the institution.", 15, INK, False, "Calibri")
add_text(s, 6.7, 2.55, 2.8, 3.2, "3. LICENCE\n\nThe purchased module pack is switched on.", 15, INK, False, "Calibri")
add_text(s, 9.7, 2.55, 2.9, 3.2, "4. USE\n\nStaff work only in the licensed portion.", 15, INK, False, "Calibri")

# 3
s = prs.slides.add_slide(blank)
kicker_title(s, "Configurable packs", "The same platform, sold as\ndifferent products.", 3)
add_text(s, 0.7, 2.55, 3.0, 2.8, "MICROFINANCE\n\nCRM, core banking,\nagency, field app", 15, INK, False, "Calibri")
add_text(s, 3.9, 2.55, 3.0, 2.8, "COOPERATIVE\n\nCRM, share register,\ncore banking", 15, INK, False, "Calibri")
add_text(s, 7.1, 2.55, 3.0, 2.8, "VSLA NETWORK\n\nCRM, share-out,\ncore banking", 15, INK, False, "Calibri")
add_text(s, 10.3, 2.55, 2.4, 2.8, "PAYGO / GAF\n\nCRM, device credit,\nOEM tokens", 15, INK, False, "Calibri")
add_text(s, 0.7, 5.5, 12, 1.2, "Packs are assembled from a catalogue and can be changed during the contract. Unlicensed modules are unavailable in the screens and in the API.", 16, MUTED, False, "Calibri")

# 4
s = prs.slides.add_slide(blank)
kicker_title(s, "Platform", "Four parts. One ledger.", 4)
add_text(s, 0.7, 2.5, 3.0, 3.4, "ODOO\n\nWhere people work: CRM, KYC, VSLA, cooperatives, PayGo, agency, System Administration.", 14, INK, False, "Calibri")
add_text(s, 3.9, 2.5, 3.0, 3.4, "APACHE FINERACT\n\nCore banking: loans, savings, schedules, balances, general ledger. APIs only.", 14, INK, False, "Calibri")
add_text(s, 7.1, 2.5, 3.0, 3.4, "FASTAPI\n\nChannel and money path. Separate process, same environment as Odoo.", 14, INK, False, "Calibri")
add_text(s, 10.3, 2.5, 2.5, 3.4, "FLUTTER\n\nField Android (offline) and customer Android. Calls FastAPI only.", 14, INK, False, "Calibri")

# 5
s = prs.slides.add_slide(blank)
kicker_title(s, "How people enter", "Staff open the operating platform.\nCustomers open the app or portal.", 5)
add_text(s, 0.7, 2.7, 5.8, 2.6, "STAFF AND ADMINISTRATORS\n\nInstitution staff and Emeraid Super Administrators work on the operating platform. The customer web portal is here too.", 16, INK, False, "Calibri")
add_text(s, 7.0, 2.7, 5.6, 2.6, "APPS AND CHANNELS\n\nField app, customer app, USSD, payment providers, and OEM callbacks use the integration service.", 16, INK, False, "Calibri")
add_text(s, 0.7, 5.5, 12, 1.2, "Core banking is not a public staff website. Money is confirmed there; people work on the operating platform and the apps.", 16, MUTED, False, "Calibri")

# 6
s = prs.slides.add_slide(blank)
kicker_title(s, "How money moves", "Every naira is posted once,\nin core banking.", 6)
add_text(s, 0.7, 2.55, 2.8, 3.0, "INSTRUCTION\n\nTeller, officer, field app, customer app, or payment", 15, INK, False, "Calibri")
add_text(s, 3.7, 2.55, 2.8, 3.0, "INTEGRATION\n\nWho is acting, licence, and branch", 15, INK, False, "Calibri")
add_text(s, 6.7, 2.55, 2.8, 3.0, "CORE BANKING\n\nPosts the transaction and updates the ledger", 15, INK, False, "Calibri")
add_text(s, 9.7, 2.55, 2.9, 3.0, "OPERATING PLATFORM\n\nShows the receipt to staff and, where appropriate, to the customer", 15, INK, False, "Calibri")
add_text(s, 0.7, 5.7, 12, 1.0, "If the network fails after posting, the same instruction is retried. A second posting is not created.", 16, MUTED, False, "Calibri")

# 7
s = prs.slides.add_slide(blank)
kicker_title(s, "Onboarding an institution", "From licence to a live institution.", 7)
add_text(s, 0.7, 2.5, 12, 4.2, "1.  Super Administrator creates the institution, branding, and module pack.\n\n2.  Head Office and branches are recorded. The institution’s administrator is invited.\n\n3.  Core banking is prepared for that institution, with a matching office for each branch.\n\n4.  Staff are given roles and branches. They sign in and see only their institution and those branches.", 16, INK, False, "Calibri")

# 8
s = prs.slides.add_slide(blank)
kicker_title(s, "Branch access", "Staff see their branch.\nMoney can only post there.", 8)
add_text(s, 0.7, 2.7, 3.8, 2.8, "OPERATING PLATFORM\n\nA Garki officer does not see Wuse customers.", 16, INK, False, "Calibri")
add_text(s, 4.8, 2.7, 3.8, 2.8, "INTEGRATION SERVICE\n\nA request for another branch is refused.", 16, INK, False, "Calibri")
add_text(s, 8.9, 2.7, 3.8, 2.8, "CORE BANKING\n\nPostings can only land in the matching branch.", 16, INK, False, "Calibri")
add_text(s, 0.7, 5.6, 12, 1.1, "Institution finance may see all branches of their own institution. They cannot see another institution.", 16, MUTED, False, "Calibri")

# 9
s = prs.slides.add_slide(blank)
kicker_title(s, "Customers", "A customer belongs to the institution,\na branch, and an officer.", 9)
add_text(s, 0.7, 2.5, 12, 4.2, "1.  Record — opened at a branch and assigned to an account officer.\n\n2.  KYC — approved on the operating platform.\n\n3.  Core banking — the same person is opened at that branch.\n\n4.  Account — savings or loan is created in core banking.\n\nThe customer application and portal show only that person’s accounts.", 16, INK, False, "Calibri")

# 10
s = prs.slides.add_slide(blank)
kicker_title(s, "Customers", "The record and the accounts are linked.\nThey are not two people.", 10)
add_text(s, 0.7, 2.7, 5.8, 3.4, "OPERATING PLATFORM\n\nName, KYC, documents, CRM, home branch, account officer. What staff use every day.", 16, INK, False, "Calibri")
add_text(s, 7.0, 2.7, 5.6, 3.4, "CORE BANKING\n\nSavings and loan accounts, balances, repayment schedule, ledger. The money.", 16, INK, False, "Calibri")
add_text(s, 0.7, 5.6, 12, 1.1, "One customer number on both sides. The officer sees their portfolio. The customer sees only themselves.", 16, MUTED, False, "Calibri")

# 11
s = prs.slides.add_slide(blank)
kicker_title(s, "Control plane", "System Administration is how\nEmeraid operates the SaaS.", 11)
add_text(s, 0.7, 2.7, 5.8, 3.2, "EMERAID SUPER ADMINISTRATOR\n\nInstitutions, licences, packs, platform connectors, preparing core banking. Multi-factor authentication and a full audit trail.", 16, INK, False, "Calibri")
add_text(s, 7.0, 2.7, 5.6, 3.2, "INSTITUTION ADMINISTRATOR\n\nTheir branches, staff, and local branding. They cannot see another institution or Emeraid’s platform credentials.", 16, INK, False, "Calibri")

# 12
s = prs.slides.add_slide(blank)
kicker_title(s, "Ownership", "Each domain has one\nsystem of record.", 12)
add_text(s, 0.7, 2.55, 12, 3.8, "Operating platform  —  institution, licence, packs, CRM, KYC, customer record, branch, account officer, VSLA, PayGo, agency.\n\nCore banking  —  savings and loan accounts, balances, general ledger.\n\nIntegration service  —  field app, customer app, USSD, payments, and OEM calls, posting money in core banking.", 18, INK, False, "Calibri")

# 13
s = prs.slides.add_slide(blank)
kicker_title(s, "Security", "Controls sit on every path,\nnot only the server.", 13)
add_text(s, 0.7, 2.6, 3.8, 3.5, "ACCESS\n\nTLS, OpenID Connect, multi-factor authentication, tenant and branch scope, licence enforcement on the API.", 15, INK, False, "Calibri")
add_text(s, 4.8, 2.6, 3.8, 3.5, "MONEY\n\nCore banking is the only ledger. Duplicate requests are not posted twice. Maker-checker. Core banking is not a public staff site.", 15, INK, False, "Calibri")
add_text(s, 8.9, 2.6, 3.8, 3.5, "OPERATIONS\n\nEncryption, NDPR, encrypted field store, audit, backup and restore, software bill of materials, penetration testing.", 15, INK, False, "Calibri")

# 14
s = prs.slides.add_slide(blank)
kicker_title(s, "Same environment", "The operating platform and the\nintegration service are neighbours.", 14)
add_text(s, 0.7, 2.55, 2.8, 3.2, "NGINX\n\nTLS at the edge", 15, INK, False, "Calibri")
add_text(s, 3.7, 2.55, 2.8, 3.2, "STAFF AND PORTAL\n\nOperating platform", 15, INK, False, "Calibri")
add_text(s, 6.7, 2.55, 2.8, 3.2, "APPS AND PARTNERS\n\nIntegration service", 15, INK, False, "Calibri")
add_text(s, 9.7, 2.55, 2.9, 3.2, "CORE BANKING\n\nInternal only", 15, INK, False, "Calibri")
add_text(s, 0.7, 5.7, 12, 1.0, "Each component keeps its own database. They connect through the integration service.", 16, MUTED, False, "Calibri")

# 15
s = prs.slides.add_slide(blank)
kicker_title(s, "End-to-end flows", "Input, process, output —\nthe same form for every process.", 15)
add_text(s, 0.7, 2.5, 3.8, 3.6, "CUSTOMER\n\nIn: details, KYC, branch, officer.\nThen: Odoo KYC; one ID; Fineract client at that branch.\nOut: approved customer, linked IDs.", 14, INK, False, "Calibri")
add_text(s, 4.8, 2.5, 3.8, 3.6, "LOAN\n\nIn: application, product, amount, tenor.\nThen: Odoo appraisal; Fineract loan at that branch.\nOut: account, schedule, audit trail.", 14, INK, False, "Calibri")
add_text(s, 8.9, 2.5, 3.8, 3.6, "REPAYMENT\n\nIn: amount, account, channel.\nThen: one posting in Fineract; Odoo receipt.\nOut: confirmed balance, no duplicate.", 14, INK, False, "Calibri")
add_text(s, 0.7, 6.15, 12, 0.6, "PayGo, VSLA share-out and offline collections use the same pattern.", 14, MUTED, False, "Calibri")

# 16
s = prs.slides.add_slide(blank)
paint_bg(s)
add_text(s, 0.7, 0.4, 12, 0.35, "SUMMARY", 13, GOLD, True, "Calibri")
add_text(s, 0.7, 1.3, 12, 2.2, "One platform.\nConfigurable packs.\nOne ledger.", 36, INK, False, "Georgia")
add_text(s, 0.7, 4.0, 11.5, 2.0, "Emeraid hosts EDFIP, onboards institutions, and licenses modules. Staff and customers use the operating platform and the apps. Core banking holds the accounts. There is one customer, one set of accounts, and one ledger.", 18, MUTED, False, "Calibri")
footer(s, 16)

out = "/workspace/edfip-screening/architecture/EDFIP_Architecture_Dexta.pptx"
prs.save(out)
print("wrote", out)
